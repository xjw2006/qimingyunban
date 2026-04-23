from __future__ import annotations

import hashlib
import json
import random
import re
import smtplib
from copy import deepcopy
from datetime import datetime, timedelta
from email.message import EmailMessage
from pathlib import Path
from typing import Any

import httpx
import requests
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PPTXColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PInches, Pt as PPt
from sqlalchemy import desc, func, select
from sqlalchemy.orm import Session

from .config import settings
from .models import BlacklistEntry, ForumComment, ForumPost, Proposal, Report, SensitiveWord, User, VerificationCode
from .security import hash_text


# ======================== 通义千问 API 配置 ========================
API_KEY = "sk-0585a7345604420cb4b18a6c84a5928e"
API_URL = "https://dashscope.aliyuncs.com/api/v1/services/aigc/text-generation/generation"

def call_cloud_llm(prompt: str) -> str:
    headers = {
        "Authorization": f"Bearer {API_KEY}",
        "Content-Type": "application/json"
    }
    data = {
        "model": "qwen-turbo",
        "input": {
            "messages": [{"role": "user", "content": prompt}]
        },
        "parameters": {
            "result_format": "text"
        }
    }
    try:
        response = requests.post(API_URL, headers=headers, json=data, timeout=60)
        response.raise_for_status()
        result = response.json()
        return result["output"]["text"]
    except Exception as e:
        return f"API调用失败：{str(e)}"
# =================================================================


EMAIL_RE = re.compile(r"^[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}$")
CHINA_PHONE_RE = re.compile(r"^1[3-9]\d{9}$")
DIGIT_RE = re.compile(r"\d")


THEME = {
    "blue": PPTXColor(14, 95, 216),
    "navy": PPTXColor(19, 49, 92),
    "light": PPTXColor(243, 247, 255),
    "sky": PPTXColor(122, 168, 255),
    "gray": PPTXColor(86, 101, 123),
    "border": PPTXColor(208, 220, 241),
}


def is_valid_email(value: str) -> bool:
    return bool(EMAIL_RE.match((value or "").strip()))


def is_valid_china_phone(value: str) -> bool:
    return bool(CHINA_PHONE_RE.match((value or "").strip()))


def is_strong_password(password: str) -> bool:
    return len(password or "") >= 8 and bool(re.search(r"[A-Za-z]", password)) and bool(re.search(r"\d", password))


def stable_number(seed_text: str, minimum: int, maximum: int) -> int:
    digest = hashlib.md5(seed_text.encode("utf-8")).hexdigest()
    span = maximum - minimum + 1
    return minimum + (int(digest[:8], 16) % span)


def send_email_verification(target: str, code: str) -> dict[str, Any]:
    preview = None
    status = "console"
    if settings.verification_delivery == "smtp" and settings.smtp_host and settings.smtp_user and settings.smtp_password:
        message = EmailMessage()
        message["From"] = settings.smtp_from_email
        message["To"] = target
        message["Subject"] = "启明云伴 邮箱验证码"
        message.set_content(
            f"您正在使用 启明云伴 商赛策划平台进行邮箱验证。\n验证码：{code}\n有效期：{settings.verification_code_expire_minutes} 分钟。"
        )
        if settings.smtp_use_ssl:
            with smtplib.SMTP_SSL(settings.smtp_host, settings.smtp_port, timeout=20) as smtp:
                smtp.login(settings.smtp_user, settings.smtp_password)
                smtp.send_message(message)
        else:
            with smtplib.SMTP(settings.smtp_host, settings.smtp_port, timeout=20) as smtp:
                smtp.starttls()
                smtp.login(settings.smtp_user, settings.smtp_password)
                smtp.send_message(message)
        status = "sent"
    else:
        preview = code if settings.dev_show_codes else None
    return {"status": status, "preview_code": preview}


def send_sms_verification(target: str, code: str) -> dict[str, Any]:
    preview = code if settings.dev_show_codes else None
    return {"status": settings.sms_provider or "console", "preview_code": preview}


def issue_verification_code(db: Session, target: str, target_type: str, purpose: str) -> dict[str, Any]:
    cooldown_start = datetime.utcnow() - timedelta(seconds=settings.verification_cooldown_seconds)
    recent_stmt = (
        select(VerificationCode)
        .where(
            VerificationCode.target == target,
            VerificationCode.target_type == target_type,
            VerificationCode.purpose == purpose,
            VerificationCode.created_at >= cooldown_start,
        )
        .order_by(desc(VerificationCode.created_at))
        .limit(1)
    )
    recent = db.execute(recent_stmt).scalar_one_or_none()
    if recent:
        wait_seconds = settings.verification_cooldown_seconds - int((datetime.utcnow() - recent.created_at).total_seconds())
        return {"ok": False, "message": f"发送过于频繁，请 {max(wait_seconds, 1)} 秒后再试。"}

    code = f"{random.randint(0, 999999):06d}"
    expires_at = datetime.utcnow() + timedelta(minutes=settings.verification_code_expire_minutes)
    record = VerificationCode(
        target=target,
        target_type=target_type,
        purpose=purpose,
        code_hash=hash_text(code),
        expires_at=expires_at,
        delivery_status="pending",
    )
    db.add(record)
    db.commit()
    db.refresh(record)

    if target_type == "email":
        send_info = send_email_verification(target, code)
    else:
        send_info = send_sms_verification(target, code)
    record.delivery_status = send_info["status"]
    db.add(record)
    db.commit()
    return {
        "ok": True,
        "message": f"验证码已发送到{mask_target(target, target_type)}。",
        "preview_code": send_info.get("preview_code"),
        "expires_in": settings.verification_code_expire_minutes * 60,
    }


def verify_code(db: Session, target: str, target_type: str, purpose: str, code: str) -> tuple[bool, str]:
    stmt = (
        select(VerificationCode)
        .where(
            VerificationCode.target == target,
            VerificationCode.target_type == target_type,
            VerificationCode.purpose == purpose,
            VerificationCode.used_at.is_(None),
        )
        .order_by(desc(VerificationCode.created_at))
        .limit(1)
    )
    record = db.execute(stmt).scalar_one_or_none()
    if not record:
        return False, "未找到验证码，请先发送。"
    if record.expires_at < datetime.utcnow():
        return False, "验证码已过期，请重新发送。"
    record.attempts += 1
    if record.attempts > 8:
        db.add(record)
        db.commit()
        return False, "尝试次数过多，请重新发送验证码。"
    if record.code_hash != hash_text(code.strip()):
        db.add(record)
        db.commit()
        return False, "验证码错误。"
    record.used_at = datetime.utcnow()
    db.add(record)
    db.commit()
    return True, "验证通过。"


def mask_target(value: str, target_type: str) -> str:
    if target_type == "phone" and len(value) == 11:
        return f"{value[:3]}****{value[-4:]}"
    if target_type == "email" and "@" in value:
        name, domain = value.split("@", 1)
        if len(name) <= 2:
            show = name[:1] + "*"
        else:
            show = name[:2] + "***"
        return f"{show}@{domain}"
    return value


DEFAULT_SENSITIVE_WORDS = [
    ("诈骗", "reject", "***"),
    ("刷单", "reject", "***"),
    ("赌博", "reject", "***"),
    ("毒品", "reject", "***"),
    ("开盒", "reject", "***"),
    ("侮辱", "replace", "不当词"),
    ("傻逼", "replace", "不当词"),
    ("傻X", "replace", "不当词"),
]


def moderate_text(db: Session, content: str) -> dict[str, Any]:
    original = content or ""
    lowered = original.lower()
    matches: list[dict[str, str]] = []
    processed = original
    reject = False

    words = db.execute(select(SensitiveWord)).scalars().all()
    blacklist_keywords = db.execute(
        select(BlacklistEntry).where(BlacklistEntry.entry_type == "keyword", BlacklistEntry.active.is_(True))
    ).scalars().all()

    for item in words:
        if item.word and item.word.lower() in lowered:
            matches.append({"source": item.word, "policy": item.level})
            if item.level == "reject":
                reject = True
            else:
                processed = re.sub(re.escape(item.word), item.replacement or "***", processed, flags=re.IGNORECASE)

    for item in blacklist_keywords:
        if item.value and item.value.lower() in lowered:
            matches.append({"source": item.value, "policy": "blacklist"})
            reject = True

    action = "allow"
    if reject:
        action = "reject"
    elif processed != original:
        action = "replace"

    return {
        "action": action,
        "content": processed,
        "matches": matches,
        "message": "内容命中敏感词或黑名单规则。" if reject else "内容已自动脱敏。" if action == "replace" else "通过。",
    }


def registration_blacklist_check(db: Session, email: str, phone: str, full_name: str) -> tuple[bool, str]:
    entries = db.execute(select(BlacklistEntry).where(BlacklistEntry.active.is_(True))).scalars().all()
    username = (full_name or "").strip().lower()
    for item in entries:
        value = (item.value or "").strip().lower()
        if item.entry_type == "email" and value == email.lower():
            return False, "该邮箱已被列入黑名单，无法注册。"
        if item.entry_type == "phone" and value == phone:
            return False, "该手机号已被列入黑名单，无法注册。"
        if item.entry_type == "username" and value == username:
            return False, "该昵称已被列入黑名单，无法注册。"
    return True, "通过。"


def seed_defaults(db: Session, password_hasher) -> None:
    if db.execute(select(func.count()).select_from(User)).scalar_one() == 0:
        demo = User(
            full_name="演示用户",
            email="demo@启明云伴.app",
            phone="13800000000",
            password_hash=password_hasher("demo1234"),
            email_verified=True,
            phone_verified=True,
            is_admin=False,
        )
        admin = User(
            full_name="平台管理员",
            email="admin@启明云伴.app",
            phone="13900000000",
            password_hash=password_hasher("admin1234"),
            email_verified=True,
            phone_verified=True,
            is_admin=True,
        )
        db.add_all([demo, admin])
        db.commit()

    if db.execute(select(func.count()).select_from(SensitiveWord)).scalar_one() == 0:
        db.add_all([SensitiveWord(word=w, level=lv, replacement=rep) for w, lv, rep in DEFAULT_SENSITIVE_WORDS])
        db.commit()

    demo_user = db.execute(select(User).where(User.email == "demo@启明云伴.app")).scalar_one_or_none()
    if demo_user and db.execute(select(func.count()).select_from(Proposal)).scalar_one() == 0:
        payload = {
            "title": "蓝鲸校园轻食连锁计划",
            "industry": "消费零售",
            "problem": "高校食堂高峰拥挤且健康轻食选择不足，学生对高蛋白、低负担的轻餐需求明显，但现有供给同质化、履约慢、体验弱。",
            "solution": "打造校园内 10 分钟即取的轻食餐柜与小程序预订体系，主打高蛋白轻食碗、定制饮品与社群会员体系。",
            "market": "首站覆盖 3 所华东高校，核心用户为 18-24 岁注重身材管理与效率的学生群体。计划通过社团合作、运动社群和宿舍团购打开首批用户。",
            "business_model": "采用餐柜零售 + 会员月卡 + 校园团购合作模式，核心收入来自单品销售、会员复购和品牌联名。",
        }
        generated = generate_proposal_payload(payload)
        score = evaluate_proposal(payload, generated)
        proposal = Proposal(
            owner_id=demo_user.id,
            title=payload["title"],
            industry=payload["industry"],
            problem=payload["problem"],
            solution=payload["solution"],
            market=payload["market"],
            business_model=payload["business_model"],
            generated_json=json.dumps(generated, ensure_ascii=False),
            score_json=json.dumps(score, ensure_ascii=False),
        )
        post = ForumPost(
            author_id=demo_user.id,
            title="如何把校园轻食项目写得更像投资人路演？",
            content="我在准备高校轻食连锁项目，想知道‘市场空间’和‘单位经济模型’两页怎么写更专业，欢迎交流。",
        )
        db.add_all([proposal, post])
        db.commit()


def maybe_generate_with_ollama(input_payload: dict[str, str]) -> dict[str, Any] | None:
    try:
        prompt = f"""
你是专业商业策划顾问，根据用户输入生成标准JSON商业计划书，字段：
summary, metric_cards(3项), sections(6项), risks(3项), judge_focus(3项), growth_roadmap(3项)
要求：专业、可用于商赛路演、纯JSON输出，不要markdown。
用户输入：{json.dumps(input_payload, ensure_ascii=False)}
""".strip()

        result = call_cloud_llm(prompt)
        start = result.find("{")
        end = result.rfind("}")
        if start == -1 or end == -1:
            return None
        return json.loads(result[start:end+1])
    except Exception:
        return None


SECTION_TITLES = [
    ("项目摘要", "一页看懂项目价值"),
    ("客群与痛点", "需求被谁触发、为何必须解决"),
    ("解决方案与差异化", "产品路径、服务链路、竞争壁垒"),
    ("市场空间与竞争格局", "市场规模、竞品结构、切入策略"),
    ("商业模式与增长飞轮", "收入来源、增长引擎、渠道设计"),
    ("落地计划与风险控制", "里程碑、关键资源、风控机制"),
]


def generate_proposal_payload(input_payload: dict[str, str]) -> dict[str, Any]:
    ollama_payload = maybe_generate_with_ollama(input_payload)
    if ollama_payload:
        return ollama_payload

    title = input_payload["title"]
    industry = input_payload["industry"]
    problem = input_payload["problem"]
    solution = input_payload["solution"]
    market = input_payload["market"]
    business_model = input_payload["business_model"]

    base_seed = f"{title}|{industry}|{problem}|{solution}|{market}|{business_model}"
    market_size = stable_number(base_seed + "market", 800, 2500)
    gross_margin = stable_number(base_seed + "margin", 42, 78)
    break_even = stable_number(base_seed + "break_even", 8, 18)
    repeat_rate = stable_number(base_seed + "repeat", 28, 74)
    acquisition_cost = stable_number(base_seed + "cac", 12, 58)

    sections = []
    section_bullets = [
        [
            f"围绕“{title}”形成一句话价值主张，强调 {industry} 场景中的效率与体验提升。",
            f"当前痛点来自：{problem[:36]}...",
            "建议首轮路演先用 30 秒讲清项目、用户、收益三件事。",
        ],
        [
            f"核心目标客群来自：{market[:34]}...",
            "用用户路径拆解“触发-决策-复购”三步，增强代入感。",
            "优先补充使用频次、决策链条、替代方案这三类证据。",
        ],
        [
            f"核心方案：{solution[:40]}...",
            "把产品链路拆成“前端触达—中端交付—后端复购”。",
            "差异化建议围绕速度、成本、体验、可复制性四条展开。",
        ],
        [
            f"目标市场规模建议首版按 {market_size} 万元级样板市场推演，再扩展到全国。",
            "竞品分析要体现：谁更强、谁更慢、谁更贵、谁更难复制。",
            "切入策略优先使用高频小场景验证，再向强关系渠道扩张。",
        ],
        [
            f"收入模型：{business_model[:36]}...",
            f"建议在路演中补充复购率 {repeat_rate}% 与获客成本 {acquisition_cost} 元的样板测算。",
            "增长飞轮建议写成：种子用户—转介绍—复购—品牌联名。",
        ],
        [
            f"盈亏平衡预计在第 {break_even} 个月附近验证，需持续追踪现金流与转化率。",
            "里程碑至少拆成 0-3 月试点、4-6 月迭代、7-12 月复制扩张。",
            "风险控制应覆盖供应、合规、履约、竞争反应四类。",
        ],
    ]
    section_narratives = [
        f"{title} 定位于 {industry} 赛道中的高可执行样板项目，核心不是简单讲创意，而是把“用户需求—解决方案—商业化路径”打通。当前版本建议先把项目价值主张压缩成一句话，再用 3 组可量化指标说明为什么值得投、值得做、值得复制。",
        f"从用户视角看，项目首先命中的是真实高频场景：{problem}。因此路演里需要把需求触发时刻、用户当前替代方案以及用户愿意为改进体验付费的理由讲清楚，才能让评委迅速理解项目的成立基础。",
        f"解决方案层面，项目不只是在卖一个产品，而是在构建完整的交付闭环：{solution}。建议展示关键环节的时间、成本与体验优势，突出差异化壁垒如何形成，并把“为什么别人难以快速复制”讲透。",
        f"市场论证建议从‘样板市场—可复制区域—全国扩张’三层展开。当前素材中已经提供了切入方向：{market}。下一步应补充市场规模假设、竞争对手分层以及渠道进入策略，让市场页面既有天花板，也有落地路径。",
        f"商业模式需要让评委看到收入从哪里来、增长靠什么跑、现金流如何稳。结合当前设定：{business_model}。建议在路演中加入会员、复购、客单价、获客成本、毛利率与回本周期等指标，构成一张完整的经营驾驶舱。",
        f"落地方面，项目需要通过分阶段目标来证明执行确定性。可以把路线图写成‘试点验证—模型稳定—渠道复制’，并针对供应链、合规、运营、竞争反应分别给出预案。这会显著提升方案的专业度与可信度。",
    ]
    chart_titles = [
        "价值逻辑图",
        "用户旅程图",
        "产品闭环图",
        "市场分层图",
        "增长飞轮图",
        "里程碑甘特图",
    ]

    for idx, (title_i, subtitle_i) in enumerate(SECTION_TITLES):
        sections.append(
            {
                "title": title_i,
                "subtitle": subtitle_i,
                "bullets": section_bullets[idx],
                "narrative": section_narratives[idx],
                "chart_placeholder": chart_titles[idx],
            }
        )

    return {
        "summary": f"{title} 面向 {industry} 场景，围绕 {problem[:20]} 的核心痛点，提供 {solution[:22]} 的解决思路，并通过 {business_model[:22]} 构建可复制的增长与盈利闭环。",
        "metric_cards": [
            {"label": "样板市场规模", "value": f"{market_size} 万元", "note": "首年聚焦样板区域"},
            {"label": "目标毛利率", "value": f"{gross_margin}%", "note": "按核心单元模型测算"},
            {"label": "盈亏平衡周期", "value": f"{break_even} 个月", "note": "随履约效率持续优化"},
        ],
        "sections": sections,
        "risks": [
            "供应链稳定性与交付质量波动",
            "竞品快速跟进导致获客成本抬升",
            "规模化扩张阶段组织效率失衡",
        ],
        "judge_focus": [
            "是否有足够明确的用户场景与支付意愿证据",
            "商业模式是否可复制、可验证、可复用",
            "增长策略能否在 6-12 个月内跑出样板数据",
        ],
        "growth_roadmap": [
            "0-3 个月：完成 MVP 与首批种子用户验证",
            "4-6 个月：打磨履约效率与复购机制",
            "7-12 个月：复制到更多高频场景与合作渠道",
        ],
        "ops_metrics": {
            "repeat_rate": f"{repeat_rate}%",
            "acquisition_cost": f"{acquisition_cost} 元",
            "sample_store_payback": f"{break_even} 个月",
        },
    }


def contains_keywords(text: str, keywords: list[str]) -> bool:
    return any(keyword in text for keyword in keywords)


def build_excerpt(text: str, length: int = 60) -> str:
    value = (text or "").strip().replace("\n", " ")
    if len(value) <= length:
        return value
    return value[:length] + "..."


SCORE_DIMENSIONS = [
    {
        "label": "痛点定义",
        "max": 15,
        "checks": [
            {"name": "明确目标客群", "points": 5, "fn": lambda s, flat, g: len(s["market"]) >= 20, "hint": "在市场描述中明确到具体用户群、场景与需求频次。", "evidence": lambda s, flat, g: build_excerpt(s["market"])},
            {"name": "痛点具备场景描述", "points": 5, "fn": lambda s, flat, g: len(s["problem"]) >= 40, "hint": "把用户使用场景、原有方案与不满意点写具体。", "evidence": lambda s, flat, g: build_excerpt(s["problem"])},
            {"name": "痛点具备量化倾向", "points": 5, "fn": lambda s, flat, g: DIGIT_RE.search(flat) is not None or contains_keywords(flat, ["高频", "复购", "%", "万元", "成本"]), "hint": "加入频次、金额、成本、时长等证据。", "evidence": lambda s, flat, g: "文案中已出现量化表达。" if (DIGIT_RE.search(flat) or contains_keywords(flat, ["%", "万元", "成本"])) else "尚缺量化数据。"},
        ],
    },
    {
        "label": "价值主张",
        "max": 15,
        "checks": [
            {"name": "解决方案对应痛点", "points": 5, "fn": lambda s, flat, g: len(s["solution"]) >= 35 and len(s["problem"]) >= 35, "hint": "让方案与痛点一一对应。", "evidence": lambda s, flat, g: build_excerpt(s["solution"])},
            {"name": "差异化路径明确", "points": 5, "fn": lambda s, flat, g: contains_keywords(flat, ["差异", "壁垒", "独特", "效率", "体验", "复制"]), "hint": "建议写出至少 2 条与竞品不同的结构性优势。", "evidence": lambda s, flat, g: "已出现差异化关键词。" if contains_keywords(flat, ["差异", "壁垒", "体验", "效率"]) else "尚未突出差异化。"},
            {"name": "结果可衡量", "points": 5, "fn": lambda s, flat, g: len(g.get("metric_cards", [])) >= 3, "hint": "加入毛利率、回本周期、转化率等结果指标。", "evidence": lambda s, flat, g: "已配置关键经营指标卡片。"},
        ],
    },
    {
        "label": "市场空间",
        "max": 15,
        "checks": [
            {"name": "市场规模思路清晰", "points": 5, "fn": lambda s, flat, g: contains_keywords(flat, ["市场", "规模", "样板", "全国", "区域"]), "hint": "建议拆分样板市场、可复制市场、全国空间。", "evidence": lambda s, flat, g: build_excerpt(s["market"])},
            {"name": "竞争格局有提及", "points": 5, "fn": lambda s, flat, g: contains_keywords(flat, ["竞品", "竞争", "替代", "同类"]), "hint": "增加竞品分层与对比矩阵。", "evidence": lambda s, flat, g: "已提及竞争/替代逻辑。" if contains_keywords(flat, ["竞品", "竞争", "替代"]) else "尚未形成竞品对比。"},
            {"name": "切入策略具体", "points": 5, "fn": lambda s, flat, g: contains_keywords(flat, ["试点", "渠道", "合作", "首批", "切入"]), "hint": "写明首批渠道、场景与扩张顺序。", "evidence": lambda s, flat, g: build_excerpt(g["growth_roadmap"][0])},
        ],
    },
    {
        "label": "商业模式",
        "max": 15,
        "checks": [
            {"name": "收入来源可解释", "points": 5, "fn": lambda s, flat, g: contains_keywords(s["business_model"], ["会员", "收费", "订阅", "佣金", "广告", "销售", "联名"]), "hint": "写清核心收入与次级收入。", "evidence": lambda s, flat, g: build_excerpt(s["business_model"])},
            {"name": "单位经济意识", "points": 5, "fn": lambda s, flat, g: contains_keywords(flat, ["毛利", "成本", "CAC", "回本", "获客"]) or len(g.get("ops_metrics", {})) >= 3, "hint": "补充 CAC、LTV、毛利率、回本周期。", "evidence": lambda s, flat, g: "已给出经营指标与样板测算。" if (contains_keywords(flat, ["毛利", "成本", "回本", "获客"]) or len(g.get("ops_metrics", {})) >= 3) else "尚缺单位经济指标。"},
            {"name": "增长飞轮成立", "points": 5, "fn": lambda s, flat, g: len(g.get("growth_roadmap", [])) >= 3, "hint": "把拉新、转化、复购、转介绍连成闭环。", "evidence": lambda s, flat, g: "已设计阶段性增长路径。"},
        ],
    },
    {
        "label": "执行计划",
        "max": 10,
        "checks": [
            {"name": "里程碑清晰", "points": 5, "fn": lambda s, flat, g: len(g.get("growth_roadmap", [])) >= 3, "hint": "至少拆成 3 个阶段。", "evidence": lambda s, flat, g: build_excerpt("；".join(g.get("growth_roadmap", [])))},
            {"name": "资源配置有方向", "points": 5, "fn": lambda s, flat, g: contains_keywords(flat, ["团队", "合作", "供应链", "运营", "履约"]), "hint": "写明关键岗位或合作资源。", "evidence": lambda s, flat, g: "已出现团队/供应链/履约等执行资源表述。" if contains_keywords(flat, ["团队", "合作", "供应链", "运营", "履约"]) else "尚未说明关键资源。"},
        ],
    },
    {
        "label": "财务逻辑",
        "max": 10,
        "checks": [
            {"name": "财务指标基础", "points": 5, "fn": lambda s, flat, g: len(g.get("metric_cards", [])) >= 3, "hint": "建议给出市场规模、毛利、回本三类指标。", "evidence": lambda s, flat, g: "已生成 3 张指标卡。"},
            {"name": "现金流/回本意识", "points": 5, "fn": lambda s, flat, g: contains_keywords(flat, ["回本", "现金流", "盈亏平衡", "毛利"]), "hint": "写出样板模型的现金流假设。", "evidence": lambda s, flat, g: "已出现回本/盈亏平衡相关表述。" if contains_keywords(flat, ["回本", "盈亏平衡", "现金流", "毛利"]) else "尚未写出回本逻辑。"},
        ],
    },
    {
        "label": "风险控制",
        "max": 10,
        "checks": [
            {"name": "风险识别齐全", "points": 5, "fn": lambda s, flat, g: len(g.get("risks", [])) >= 3, "hint": "建议至少覆盖供应、合规、竞争三类风险。", "evidence": lambda s, flat, g: build_excerpt("；".join(g.get("risks", [])))},
            {"name": "应对思路存在", "points": 5, "fn": lambda s, flat, g: contains_keywords(flat, ["预案", "控制", "应对", "优化", "机制"]), "hint": "用一页写清每类风险的缓释动作。", "evidence": lambda s, flat, g: "已出现风控/预案类表述。" if contains_keywords(flat, ["预案", "控制", "应对", "机制"]) else "尚缺具体应对动作。"},
        ],
    },
    {
        "label": "展示完成度",
        "max": 10,
        "checks": [
            {"name": "章节结构完整", "points": 5, "fn": lambda s, flat, g: len(g.get("sections", [])) >= 6, "hint": "建议围绕 6 大核心章节展开。", "evidence": lambda s, flat, g: f"已生成 {len(g.get('sections', []))} 个核心章节。"},
            {"name": "评委关注点明确", "points": 5, "fn": lambda s, flat, g: len(g.get("judge_focus", [])) >= 3, "hint": "单独写出评委最容易追问的 3 个问题。", "evidence": lambda s, flat, g: build_excerpt("；".join(g.get("judge_focus", [])))},
        ],
    },
]


def evaluate_proposal(input_payload: dict[str, str], generated: dict[str, Any]) -> dict[str, Any]:
    flat_text = " ".join(
        [
            input_payload.get("title", ""),
            input_payload.get("industry", ""),
            input_payload.get("problem", ""),
            input_payload.get("solution", ""),
            input_payload.get("market", ""),
            input_payload.get("business_model", ""),
            generated.get("summary", ""),
            " ".join(section.get("narrative", "") for section in generated.get("sections", [])),
            " ".join(item.get("value", "") for item in generated.get("metric_cards", [])),
        ]
    )
    dimensions: list[dict[str, Any]] = []
    total = 0
    max_total = 0
    for dimension in SCORE_DIMENSIONS:
        subtotal = 0
        criteria = []
        for item in dimension["checks"]:
            passed = bool(item["fn"](input_payload, flat_text, generated))
            earned = item["points"] if passed else 0
            subtotal += earned
            criteria.append(
                {
                    "name": item["name"],
                    "points": item["points"],
                    "earned": earned,
                    "passed": passed,
                    "evidence": item["evidence"](input_payload, flat_text, generated),
                    "how_to_gain": item["hint"],
                }
            )
        total += subtotal
        max_total += dimension["max"]
        dimensions.append(
            {
                "label": dimension["label"],
                "score": subtotal,
                "max": dimension["max"],
                "criteria": criteria,
                "score_gap": dimension["max"] - subtotal,
                "comment": "该维度已具备较好说服力。" if subtotal >= dimension["max"] * 0.8 else "该维度还有明显可提升空间。",
            }
        )
    if total >= 88:
        grade = "A+"
    elif total >= 80:
        grade = "A"
    elif total >= 70:
        grade = "B+"
    elif total >= 60:
        grade = "B"
    else:
        grade = "C"
    improve_first = sorted(dimensions, key=lambda item: item["score_gap"], reverse=True)[:3]
    return {
        "total": total,
        "max_total": max_total,
        "grade": grade,
        "dimensions": dimensions,
        "summary": f"当前总分 {total}/{max_total}，评级 {grade}。此评分采用加分制，每一分均由可解释的结构化项累加而来。",
        "priority_actions": [
            f"优先补强“{item['label']}”维度，可新增 {item['score_gap']} 分。" for item in improve_first if item["score_gap"] > 0
        ],
    }


def _set_doc_table_borders(table):
    tbl = table._tbl
    tbl_pr = tbl.tblPr
    borders = OxmlElement("w:tblBorders")
    for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
        elem = OxmlElement(f"w:{edge}")
        elem.set(qn("w:val"), "single")
        elem.set(qn("w:sz"), "8")
        elem.set(qn("w:space"), "0")
        elem.set(qn("w:color"), "D0DCF1")
        borders.append(elem)
    tbl_pr.append(borders)


def export_proposal_docx(proposal: Proposal, generated: dict[str, Any], score: dict[str, Any], owner_name: str) -> Path:
    document = Document()
    section = document.sections[0]
    section.top_margin = Inches(0.7)
    section.bottom_margin = Inches(0.6)
    section.left_margin = Inches(0.8)
    section.right_margin = Inches(0.8)

    title = document.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run(proposal.title)
    run.bold = True
    run.font.size = Pt(24)
    run.font.color.rgb = RGBColor(14, 95, 216)

    sub = document.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub_run = sub.add_run(f"{settings.brand_name} · 商业策划书 V2.3 · 生成时间 {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}")
    sub_run.font.size = Pt(10)
    sub_run.font.color.rgb = RGBColor(86, 101, 123)

    summary_head = document.add_paragraph()
    summary_head.style = "Heading 1"
    summary_head.add_run("项目摘要")
    paragraph = document.add_paragraph(generated["summary"])
    paragraph.style = document.styles["Normal"]

    metrics_title = document.add_paragraph()
    metrics_title.style = "Heading 1"
    metrics_title.add_run("关键指标")
    metrics_table = document.add_table(rows=1, cols=3)
    metrics_table.style = "Table Grid"
    _set_doc_table_borders(metrics_table)
    hdr = metrics_table.rows[0].cells
    hdr[0].text = "指标"
    hdr[1].text = "数值"
    hdr[2].text = "说明"
    for card in generated.get("metric_cards", []):
        row = metrics_table.add_row().cells
        row[0].text = card.get("label", "")
        row[1].text = card.get("value", "")
        row[2].text = card.get("note", "")

    for section_data in generated.get("sections", []):
        heading = document.add_paragraph()
        heading.style = "Heading 1"
        heading.add_run(section_data["title"])
        subtitle = document.add_paragraph(section_data.get("subtitle", ""))
        subtitle.runs[0].italic = True if subtitle.runs else False
        for bullet in section_data.get("bullets", []):
            document.add_paragraph(bullet, style="List Bullet")
        document.add_paragraph(section_data.get("narrative", ""))

    score_heading = document.add_paragraph()
    score_heading.style = "Heading 1"
    score_heading.add_run("加分制评分拆解")
    score_para = document.add_paragraph(score.get("summary", ""))
    score_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
    for dimension in score.get("dimensions", []):
        document.add_paragraph(f"{dimension['label']}：{dimension['score']}/{dimension['max']}", style="Heading 2")
        score_table = document.add_table(rows=1, cols=4)
        score_table.style = "Table Grid"
        _set_doc_table_borders(score_table)
        cells = score_table.rows[0].cells
        cells[0].text = "评分项"
        cells[1].text = "得分"
        cells[2].text = "证据"
        cells[3].text = "如何加分"
        for item in dimension.get("criteria", []):
            row = score_table.add_row().cells
            row[0].text = item["name"]
            row[1].text = f"+{item['earned']} / {item['points']}"
            row[2].text = item["evidence"]
            row[3].text = item["how_to_gain"]
        document.add_paragraph(dimension.get("comment", ""))

    footer = document.add_paragraph()
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer_run = footer.add_run(f"版权所有 © {settings.brand_name} · 导出用户：{owner_name}")
    footer_run.font.size = Pt(9)
    footer_run.font.color.rgb = RGBColor(86, 101, 123)

    file_path = settings.export_path / f"proposal-{proposal.id}.docx"
    document.save(file_path)
    return file_path


def _add_footer(slide, slide_no: int):
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, PInches(0), PInches(7.1), PInches(13.33), PInches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = THEME["blue"]
    line.line.fill.background()

    footer = slide.shapes.add_textbox(PInches(0.5), PInches(7.18), PInches(5.5), PInches(0.2))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = settings.brand_name
    p.font.size = PPt(9)
    p.font.color.rgb = THEME["gray"]

    page = slide.shapes.add_textbox(PInches(12.0), PInches(7.15), PInches(0.8), PInches(0.2))
    pt = page.text_frame.paragraphs[0]
    pt.text = str(slide_no)
    pt.alignment = PP_ALIGN.RIGHT
    pt.font.size = PPt(9)
    pt.font.color.rgb = THEME["gray"]


def _add_title_block(slide, title: str, subtitle: str | None = None):
    title_box = slide.shapes.add_textbox(PInches(0.7), PInches(0.45), PInches(7.6), PInches(0.8))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = PPt(24)
    p.font.bold = True
    p.font.color.rgb = THEME["navy"]
    if subtitle:
        sub_box = slide.shapes.add_textbox(PInches(0.72), PInches(1.0), PInches(7.5), PInches(0.4))
        sp = sub_box.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = PPt(11)
        sp.font.color.rgb = THEME["gray"]


def _add_metric_card(slide, x: float, y: float, title: str, value: str, note: str):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PInches(x), PInches(y), PInches(3.8), PInches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = THEME["light"]
    box.line.color.rgb = THEME["border"]
    t = slide.shapes.add_textbox(PInches(x + 0.2), PInches(y + 0.12), PInches(3.4), PInches(0.25))
    tp = t.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = PPt(10)
    tp.font.color.rgb = THEME["gray"]
    v = slide.shapes.add_textbox(PInches(x + 0.2), PInches(y + 0.35), PInches(3.4), PInches(0.38))
    vp = v.text_frame.paragraphs[0]
    vp.text = value
    vp.font.size = PPt(22)
    vp.font.bold = True
    vp.font.color.rgb = THEME["blue"]
    n = slide.shapes.add_textbox(PInches(x + 0.2), PInches(y + 0.8), PInches(3.2), PInches(0.25))
    np = n.text_frame.paragraphs[0]
    np.text = note
    np.font.size = PPt(9)
    np.font.color.rgb = THEME["gray"]


def _add_content_panel(slide, section_title: str, bullets: list[str], narrative: str, placeholder: str, slide_no: int):
    _add_title_block(slide, section_title)
    left_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PInches(0.7), PInches(1.45), PInches(5.6), PInches(5.15))
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = THEME["light"]
    left_panel.line.color.rgb = THEME["border"]

    content = slide.shapes.add_textbox(PInches(1.0), PInches(1.7), PInches(4.9), PInches(1.6))
    tf = content.text_frame
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = PPt(14)
        p.font.color.rgb = THEME["navy"]
        p.space_after = PPt(6)
    narrative_box = slide.shapes.add_textbox(PInches(1.0), PInches(3.35), PInches(4.9), PInches(2.8))
    np = narrative_box.text_frame.paragraphs[0]
    np.text = narrative
    np.font.size = PPt(12)
    np.font.color.rgb = THEME["gray"]

    right_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PInches(6.55), PInches(1.45), PInches(6.0), PInches(5.15))
    right_panel.fill.solid()
    right_panel.fill.fore_color.rgb = PPTXColor(255, 255, 255)
    right_panel.line.color.rgb = THEME["border"]
    header_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, PInches(6.55), PInches(1.45), PInches(6.0), PInches(0.45))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = THEME["blue"]
    header_bar.line.fill.background()

    hp = slide.shapes.add_textbox(PInches(6.8), PInches(1.56), PInches(5.3), PInches(0.2)).text_frame.paragraphs[0]
    hp.text = f"图表占位 · {placeholder}"
    hp.font.size = PPt(11)
    hp.font.bold = True
    hp.font.color.rgb = PPTXColor(255, 255, 255)

    chart_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PInches(6.95), PInches(2.1), PInches(5.2), PInches(3.35))
    chart_box.fill.solid()
    chart_box.fill.fore_color.rgb = THEME["light"]
    chart_box.line.color.rgb = THEME["sky"]
    chart_text = slide.shapes.add_textbox(PInches(7.25), PInches(3.1), PInches(4.6), PInches(0.8))
    ctp = chart_text.text_frame.paragraphs[0]
    ctp.text = "建议放置：样板市场规模图、用户增长曲线\n或单位经济模型对比图"
    ctp.alignment = PP_ALIGN.CENTER
    ctp.font.size = PPt(14)
    ctp.font.color.rgb = THEME["gray"]

    _add_footer(slide, slide_no)


def export_proposal_pptx(proposal: Proposal, generated: dict[str, Any], score: dict[str, Any]) -> Path:
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, PInches(0), PInches(0), PInches(13.333), PInches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PPTXColor(247, 250, 255)
    bg.line.fill.background()
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, PInches(0), PInches(0), PInches(13.333), PInches(1.2))
    band.fill.solid()
    band.fill.fore_color.rgb = THEME["blue"]
    band.line.fill.background()
    title_box = slide.shapes.add_textbox(PInches(0.8), PInches(1.55), PInches(8.7), PInches(1.4))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = proposal.title
    tp.font.size = PPt(28)
    tp.font.bold = True
    tp.font.color.rgb = THEME["navy"]
    sub_box = slide.shapes.add_textbox(PInches(0.82), PInches(2.5), PInches(9.2), PInches(0.8))
    sp = sub_box.text_frame.paragraphs[0]
    sp.text = f"{proposal.industry} · 商业策划书路演版 · {settings.brand_name}"
    sp.font.size = PPt(14)
    sp.font.color.rgb = THEME["gray"]
    summary_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PInches(0.82), PInches(3.3), PInches(7.2), PInches(1.75))
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = THEME["light"]
    summary_box.line.color.rgb = THEME["border"]
    sb = slide.shapes.add_textbox(PInches(1.08), PInches(3.55), PInches(6.6), PInches(1.25))
    sbp = sb.text_frame.paragraphs[0]
    sbp.text = generated.get("summary", "")
    sbp.font.size = PPt(13)
    sbp.font.color.rgb = THEME["gray"]
    for idx, card in enumerate(generated.get("metric_cards", [])[:3]):
        _add_metric_card(slide, 8.4, 1.9 + idx * 1.55, card.get("label", ""), card.get("value", ""), card.get("note", ""))
    _add_footer(slide, 1)

    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, PInches(0), PInches(0), PInches(13.333), PInches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = THEME["blue"]
    bg.line.fill.background()
    title_b = slide.shapes.add_textbox(PInches(0.9), PInches(2.0), PInches(8.0), PInches(0.8))
    p = title_b.text_frame.paragraphs[0]
    p.text = "01｜商业逻辑与项目价值"
    p.font.size = PPt(30)
    p.font.bold = True
    p.font.color.rgb = PPTXColor(255, 255, 255)
    sub_b = slide.shapes.add_textbox(PInches(0.92), PInches(2.9), PInches(8.8), PInches(0.6))
    sp = sub_b.text_frame.paragraphs[0]
    sp.text = "从项目摘要、用户痛点到产品差异化，形成完整的价值证明链。"
    sp.font.size = PPt(14)
    sp.font.color.rgb = PPTXColor(230, 240, 255)
    _add_footer(slide, 2)

    slide_no = 3
    for idx, section_data in enumerate(generated.get("sections", []), start=1):
        slide = prs.slides.add_slide(blank)
        bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, PInches(0), PInches(0), PInches(13.333), PInches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = PPTXColor(255, 255, 255)
        bg.line.fill.background()
        _add_content_panel(
            slide,
            f"{idx:02d}｜{section_data.get('title', '')}",
            section_data.get("bullets", []),
            section_data.get("narrative", ""),
            section_data.get("chart_placeholder", "图表占位"),
            slide_no,
        )
        slide_no += 1

    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, PInches(0), PInches(0), PInches(13.333), PInches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PPTXColor(255, 255, 255)
    bg.line.fill.background()
    _add_title_block(slide, "评分拆解页", f"当前总分 {score['total']}/{score['max_total']} · 评级 {score['grade']} · 采用可解释加分制")
    y = 1.55
    x_positions = [0.7, 4.45, 8.2]
    for idx, dimension in enumerate(score.get("dimensions", [])):
        x = x_positions[idx % 3]
        if idx and idx % 3 == 0:
            y += 2.05
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PInches(x), PInches(y), PInches(3.35), PInches(1.75))
        box.fill.solid()
        box.fill.fore_color.rgb = THEME["light"]
        box.line.color.rgb = THEME["border"]
        t1 = slide.shapes.add_textbox(PInches(x + 0.18), PInches(y + 0.12), PInches(3.0), PInches(0.25))
        p1 = t1.text_frame.paragraphs[0]
        p1.text = dimension["label"]
        p1.font.size = PPt(11)
        p1.font.color.rgb = THEME["gray"]
        t2 = slide.shapes.add_textbox(PInches(x + 0.18), PInches(y + 0.45), PInches(2.8), PInches(0.35))
        p2 = t2.text_frame.paragraphs[0]
        p2.text = f"{dimension['score']} / {dimension['max']}"
        p2.font.size = PPt(22)
        p2.font.bold = True
        p2.font.color.rgb = THEME["blue"]
        t3 = slide.shapes.add_textbox(PInches(x + 0.18), PInches(y + 0.93), PInches(3.0), PInches(0.6))
        p3 = t3.text_frame.paragraphs[0]
        p3.text = dimension["comment"]
        p3.font.size = PPt(9)
        p3.font.color.rgb = THEME["gray"]
    _add_footer(slide, slide_no)
    slide_no += 1

    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, PInches(0), PInches(0), PInches(13.333), PInches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = THEME["navy"]
    bg.line.fill.background()
    title = slide.shapes.add_textbox(PInches(0.9), PInches(2.0), PInches(9.0), PInches(0.8))
    tp = title.text_frame.paragraphs[0]
    tp.text = "谢谢评审 · 进入问答环节"
    tp.font.size = PPt(30)
    tp.font.bold = True
    tp.font.color.rgb = PPTXColor(255, 255, 255)
    bullet_box = slide.shapes.add_textbox(PInches(0.95), PInches(3.0), PInches(8.2), PInches(1.5))
    tf = bullet_box.text_frame
    for idx, text in enumerate(score.get("priority_actions", [])[:3] or ["可继续补充实证数据、竞品对比与单位经济测算。"]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {text}"
        p.font.size = PPt(15)
        p.font.color.rgb = PPTXColor(230, 240, 255)
    brand = slide.shapes.add_textbox(PInches(0.95), PInches(6.2), PInches(6.5), PInches(0.3))
    bp = brand.text_frame.paragraphs[0]
    bp.text = settings.brand_name
    bp.font.size = PPt(12)
    bp.font.color.rgb = PPTXColor(187, 209, 255)
    _add_footer(slide, slide_no)

    file_path = settings.export_path / f"proposal-{proposal.id}.pptx"
    prs.save(file_path)
    return file_path


def proposal_to_dict(proposal: Proposal) -> dict[str, Any]:
    return {
        "id": proposal.id,
        "title": proposal.title,
        "industry": proposal.industry,
        "problem": proposal.problem,
        "solution": proposal.solution,
        "market": proposal.market,
        "business_model": proposal.business_model,
        "generated": json.loads(proposal.generated_json),
        "score": json.loads(proposal.score_json),
    }


def serialize_posts_for_ui(posts: list[ForumPost], current_user: User | None) -> list[dict[str, Any]]:
    post_items = []
    for post in posts:
        visible_comments = []
        for comment in sorted(post.comments, key=lambda x: x.created_at):
            if comment.is_hidden and not (current_user and current_user.is_admin):
                continue
            visible_comments.append(comment)
        if post.is_hidden and not (current_user and current_user.is_admin):
            continue
        post_items.append({"post": post, "comments": visible_comments})
    return post_items


def report_summary(db: Session) -> dict[str, int]:
    rows = db.execute(select(Report.status, func.count()).group_by(Report.status)).all()
    return {status: count for status, count in rows}